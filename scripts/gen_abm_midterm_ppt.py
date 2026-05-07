"""ABM 중간발표 PPT — 시각 중심 (텍스트 최소화).

8장. 슬라이드당 핵심 3~4 불릿, 큰 숫자/모듈명 hero, 충분한 여백.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "presentation" / "abm-midterm-2026-05-07.pptx"

NAVY = RGBColor(0x1F, 0x2D, 0x4E)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GRAY = RGBColor(0x37, 0x41, 0x51)
SUB = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)


def add_title(slide, num: int, text: str) -> None:
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.5), Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    btf = badge.text_frame
    btf.margin_left = Inches(0)
    btf.margin_right = Inches(0)
    btf.margin_top = Inches(0)
    btf.margin_bottom = Inches(0)
    bp = btf.paragraphs[0]
    bp.text = str(num)
    bp.alignment = 2
    bp.font.size = Pt(22)
    bp.font.bold = True
    bp.font.color.rgb = WHITE

    tb = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(11), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_bullets(slide, items: list[str], left=1.0, top=2.0, width=11.3, size=22) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY
        p.space_after = Pt(18)


def hero_number(slide, left: float, top: float, big: str, label: str, color=ACCENT, w: float = 3.0) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = big
    p.alignment = 2
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = label
    p2.alignment = 2
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY


def chip(slide, left: float, top: float, text: str, color=NAVY, w: float = 3.5) -> None:
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.45))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    tf = c.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.name = "Consolas"
    p.font.bold = True
    p.font.color.rgb = WHITE


# ============================================================
# 1 — 문제
# ============================================================
def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 1, "문제 — 비용 vs 다양성")

    hero_number(slide, 1.0, 2.0, "$0.37", "v1 run당 LLM 비용", RED, 3.5)
    hero_number(slide, 4.9, 2.0, "0.58~0.76", "v1 Pearson r", ORANGE, 3.5)
    hero_number(slide, 8.8, 2.0, "1000+", "마포 에이전트", ACCENT, 3.5)

    add_bullets(
        slide,
        [
            "매 tick LLM 호출 → agent·일수 늘릴수록 비용 폭증",
            "룰 기반만 쓰면 시간대·날씨 반응 다양성 손실",
            "Optuna 100 trials = 비용 폭발 → 튜닝 불가",
        ],
        top=4.5,
        size=20,
    )


# ============================================================
# 2 — 아키텍처
# ============================================================
def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 2, "3 Layer 하이브리드")

    layers = [
        ("L1", "Policy\nGeneration", "LLM 11회", ACCENT),
        ("L2", "Time +\nArchetype", "곱셈 테이블", GREEN),
        ("L3", "Agent\nDecision", "Tier S/A/B", ORANGE),
    ]
    for i, (tag, title, sub, color) in enumerate(layers):
        x = 1.0 + i * 4.15
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(3.7), Inches(3.3))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.4)
        p = tf.paragraphs[0]
        p.text = tag
        p.alignment = 2
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = title
        p2.alignment = 2
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(20)
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.alignment = 2
        p3.font.size = Pt(16)
        p3.font.color.rgb = LIGHT
        p3.space_before = Pt(12)

        if i < 2:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 3.75), Inches(3.6), Inches(0.4), Inches(0.6))
            arr.fill.solid()
            arr.fill.fore_color.rgb = NAVY
            arr.line.fill.background()

    chip(slide, 1.0, 6.2, "policy_generator.py", ACCENT, w=3.7)
    chip(slide, 5.15, 6.2, "archetypes.py", GREEN, w=3.7)
    chip(slide, 9.3, 6.2, "brain.py + policy_executor.py", ORANGE, w=3.7)


# ============================================================
# 3 — Layer 1
# ============================================================
def slide_layer1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 3, "Layer 1 — Policy Generation")

    hero_number(slide, 1.0, 2.0, "11", "LLM 호출 (1회)", ACCENT, 3.0)
    hero_number(slide, 4.4, 2.0, "6 × 2", "role × weather", GREEN, 3.0)
    hero_number(slide, 7.8, 2.0, "0", "재실행 시 호출", PURPLE, 3.0)

    add_bullets(
        slide,
        [
            "역할 6종 × 날씨 2종 = 11 base PersonaPolicy",
            "PersonaPolicy = 0~1 float 12 필드 + dong_affinity",
            "OpenAI gpt-4o-mini → Ollama qwen2.5:3b → mock",
            "policy_cache.json 영속 → 재실행 0회",
        ],
        top=4.5,
        size=20,
    )


# ============================================================
# 4 — Layer 2
# ============================================================
def slide_layer2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 4, "Layer 2 — 시간 확장 + Archetype")

    hero_number(slide, 1.0, 2.0, "11 → 66", "정책 확장", GREEN, 3.5)
    hero_number(slide, 4.9, 2.0, "30+", "Archetype 종류", ACCENT, 3.5)
    hero_number(slide, 8.8, 2.0, "±15%", "개체 jitter", ORANGE, 3.5)

    add_bullets(
        slide,
        [
            "_TIME_BLOCK_DELTAS — morning/lunch/afternoon/evening/night 곱셈",
            "_ROLE_TIME_OVERRIDES — ext_visitor evening pub ×1.7 등",
            "Archetype: homebody / night_owl / trendy_local / fitness ...",
            "1000명 × 66 정책 × ±15% = 사실상 1000 고유 행동",
        ],
        top=4.5,
        size=20,
    )


# ============================================================
# 5 — Layer 3
# ============================================================
def slide_layer3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 5, "Layer 3 — Tier 라우팅")

    tiers = [
        ("S", "50명", "Haiku 4.5\n+ ephemeral cache", "풀 LLM + 서사", PURPLE),
        ("A", "200명", "Gemini 2.5\nFlash-Lite", "경량 LLM", GREEN),
        ("B", "750명", "policy_executor\n순수 Python", "결정적 함수", ORANGE),
    ]
    for i, (tag, n, model, role, color) in enumerate(tiers):
        x = 1.0 + i * 4.15
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.7), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = color
        card.line.width = Pt(3)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = f"Tier {tag}"
        p.alignment = 2
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = n
        p2.alignment = 2
        p2.font.size = Pt(36)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_before = Pt(8)
        p3 = tf.add_paragraph()
        p3.text = model
        p3.alignment = 2
        p3.font.size = Pt(15)
        p3.font.color.rgb = GRAY
        p3.space_before = Pt(20)
        p4 = tf.add_paragraph()
        p4.text = role
        p4.alignment = 2
        p4.font.size = Pt(13)
        p4.font.italic = True
        p4.font.color.rgb = SUB
        p4.space_before = Pt(8)

    box = slide.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = "Auto-downgrade: anthropic → openai → ollama → mock"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = SUB
    p.alignment = 1


# ============================================================
# 6 — 메모리·상태·소셜
# ============================================================
def slide_memory(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 6, "기억 · 상태 · 소셜")

    cards = [
        ("Layer 2", "기억", ["visit_history", "learned_prefs", "blacklist", "habit_store"], ACCENT),
        ("Layer 3", "상태", ["hunger", "fatigue", "mood", "tick decay/recover"], GREEN),
        ("Layer 5", "소셜", ["만족도 > 0.7", "친구 2명 전파", "다음 tick visit_p ↑", "매장 평판 확산"], ORANGE),
    ]
    for i, (lab, title, items, color) in enumerate(cards):
        x = 1.0 + i * 4.15
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(3.7), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2.5)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = lab
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        p2.space_after = Pt(15)
        for it in items:
            pi = tf.add_paragraph()
            pi.text = f"· {it}"
            pi.font.size = Pt(16)
            pi.font.color.rgb = GRAY
            pi.space_after = Pt(6)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "Memory Seeder"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p2 = tf.add_paragraph()
    p2.text = "14일 가상 visit 주입 → cold start 완화"
    p2.font.size = Pt(18)
    p2.font.color.rgb = GRAY


# ============================================================
# 7 — 동 DNA
# ============================================================
def slide_dong(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 7, "마포 16동 상권 DNA")

    types = [
        ("nightlife", "서교 · 합정 · 신수", "주점 1.5", PURPLE),
        ("trendy", "연남 · 망원1", "카페 1.5", RGBColor(0xEC, 0x48, 0x99)),
        ("office", "상암 · 공덕 · 도화", "음식점 1.3", ACCENT),
        ("residential", "용강 · 아현 · 염리 · 대흥 · 서강 · 성산", "편의점 1.2", GREEN),
        ("traditional", "망원2", "혼합 1.2", ORANGE),
    ]
    for i, (tname, dongs, boost, color) in enumerate(types):
        y = 2.0 + i * 0.9
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y), Inches(2.2), Inches(0.7))
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        tf = tag.text_frame
        tf.margin_top = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = tname
        p.alignment = 2
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        db = slide.shapes.add_textbox(Inches(3.4), Inches(y + 0.05), Inches(7.0), Inches(0.6))
        p = db.text_frame.paragraphs[0]
        p.text = dongs
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY

        bb = slide.shapes.add_textbox(Inches(10.5), Inches(y + 0.05), Inches(2.5), Inches(0.6))
        p = bb.text_frame.paragraphs[0]
        p.text = boost
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = 2

    box = slide.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = "DONG_CHARACTER → 매장 score에 cat_boost 직접 곱"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = SUB
    p.alignment = 1


# ============================================================
# 8 — 결과
# ============================================================
def slide_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 8, "결과 — 학술 최첨단 수준 + 비용 $0")

    hero_number(slide, 1.0, 1.9, "4.6%", "RMSE", GREEN, 3.5)
    hero_number(slide, 4.9, 1.9, "0.69", "Pearson r", ACCENT, 3.5)
    hero_number(slide, 8.8, 1.9, "92%", "External 귀환", PURPLE, 3.5)

    # 학술 비교 테이블
    rows = [
        ["연구 / 모델", "Pearson r", "RMSE / MAPE", "단위"],
        ["우리 ABM (v12)", "0.688", "4.6%", "visit ↔ presence"],
        ["Crols & Malleson 2019", "0.7~0.9", "MAPE 10~25%", "pedestrian count (최첨단)"],
        ["Sommet & Lipps 2025", "0.5~0.85", "—", "panel stock"],
        ["Brussels 2024 (천장)", "0.96", "—", "trip ↔ trip (같은 단위)"],
    ]
    table_shape = slide.shapes.add_table(len(rows), 4, Inches(1.0), Inches(4.1), Inches(11.3), Inches(2.4))
    tbl = table_shape.table
    widths = [Inches(3.5), Inches(2.0), Inches(2.5), Inches(3.3)]
    for i, w in enumerate(widths):
        tbl.columns[i].width = w
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.alignment = 1
                for run in para.runs:
                    run.font.size = Pt(15)
                    run.font.color.rgb = WHITE if r == 0 else GRAY
                    run.font.bold = r == 0 or "우리" in row_data[0]
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif "우리" in row_data[0]:
                cell.fill.fore_color.rgb = RGBColor(0xDC, 0xFC, 0xE7)
            else:
                cell.fill.fore_color.rgb = LIGHT if r % 2 == 0 else WHITE

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.65), Inches(11.3), Inches(0.55))
    note.fill.solid()
    note.fill.fore_color.rgb = LIGHT
    note.line.color.rgb = GREEN
    note.line.width = Pt(1.5)
    tf = note.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = "v12 = Crols 2019 최첨단(Pearson 0.7~0.9) 하단 + RMSE는 그보다 낮음 — 비용 $0로 달성"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY


# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_problem(prs)
    slide_architecture(prs)
    slide_layer1(prs)
    slide_layer2(prs)
    slide_layer3(prs)
    slide_memory(prs)
    slide_dong(prs)
    slide_results(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"saved: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
